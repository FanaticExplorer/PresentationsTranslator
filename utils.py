import os
from pptx import Presentation
from googletrans import Translator
from concurrent.futures import ThreadPoolExecutor, as_completed


def translate_text(text, dest_language="uk"):
    translator = Translator()
    try:
        translated = translator.translate(text, dest=dest_language)
        return translated.text
    except Exception as e:
        print(f"Translation error: {e}")
        return text


def add_spaces_around(run, previous_run, next_run):
    """Adds spaces before and after the text in a run if it's in the middle of other text."""
    text = run.text

    # Add space before if there's a previous run and the current text doesn't start with a space
    if (
        previous_run
        and not previous_run.text.endswith(" ")
        and not text.startswith(" ")
    ):
        text = " " + text

    # Add space after if there's a next run and the current text doesn't end with a space
    if next_run and not next_run.text.startswith(" ") and not text.endswith(" "):
        text = text + " "

    run.text = text


def translate_shape(shape, dest_language="uk"):
    if shape.has_text_frame:
        paragraphs = shape.text_frame.paragraphs
        for paragraph in paragraphs:
            for i, run in enumerate(paragraph.runs):
                previous_run = paragraph.runs[i - 1] if i > 0 else None
                next_run = (
                    paragraph.runs[i + 1] if i < len(paragraph.runs) - 1 else None
                )

                original_text = run.text
                translated_text = translate_text(original_text, dest_language)

                # Apply translation and fix spaces if necessary
                run.text = translated_text
                add_spaces_around(run, previous_run, next_run)


def translate_slide(slide, dest_language="uk"):
    for shape in slide.shapes:
        translate_shape(shape, dest_language)


def translate_powerpoint(input_file, output_file, dest_language="uk", max_workers=10):
    prs = Presentation(input_file)

    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        futures = {
            executor.submit(translate_slide, slide, dest_language): slide
            for slide in prs.slides
        }
        for future in as_completed(futures):
            try:
                future.result()
            except Exception as e:
                print(f"Error processing slide: {e}")

    prs.save(output_file)


if __name__ == "__main__":
    input_file = input("Enter the path to the input PowerPoint file: ")
    output_file = os.path.splitext(input_file)[0] + "_ukrainian.pptx"

    translate_powerpoint("test.pptx", "test_ukrainian.pptx")
    print(f"Translated PowerPoint saved as: {output_file}")
